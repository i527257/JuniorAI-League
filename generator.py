import sys
import json
import os
import traceback
import urllib.request  

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    def create_slide(prs, slide_info, stijl, output_dir):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        layout = slide_info.get('layout', 'content')
        img_url = slide_info.get('img_url', '') # Pak de kant-en-klare Unsplash link van Node.js
        
        has_image = (img_url != "" and layout != "title")

        styles = {
            "futuristic": {"bg": (12, 14, 36), "card": (20, 24, 56), "text": (0, 255, 204), "accent": (255, 0, 128)},
            "playful": {"bg": (240, 248, 255), "text": (255, 255, 255), "accent": (255, 94, 0), "accent2": (255, 238, 173)},
            "professional": {"bg": (244, 245, 247), "text": (18, 30, 49), "accent": (15, 32, 67), "accent2": (0, 100, 255)},
            "modern": {"bg": (255, 255, 255), "text": (20, 20, 20), "accent": (94, 53, 177), "accent2": (235, 225, 255)}
        }
        config = styles.get(stijl, styles["modern"])
        
        # ==========================================
        # STAP 1: ACHTERGROND ONTWERP
        # ==========================================
        bg_shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
        bg_shape.fill.solid()
        bg_shape.line.fill.background()

        if stijl == "futuristic":
            bg_shape.fill.fore_color.rgb = RGBColor(*config["bg"])
            if layout == "title":
                frame = slide.shapes.add_shape(1, Inches(1.5), Inches(1.2), Inches(7.0), Inches(3.2))
            else:
                sidebar = slide.shapes.add_shape(1, 0, 0, Inches(0.3), Inches(5.625))
                sidebar.fill.solid()
                sidebar.fill.fore_color.rgb = RGBColor(*config["accent"])
                sidebar.line.fill.background()
                c_width = Inches(4.5) if has_image else Inches(8.8)
                frame = slide.shapes.add_shape(1, Inches(0.6), Inches(1.5), c_width, Inches(3.5))
            
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(*config["card"])
            frame.line.color.rgb = RGBColor(*config["text"])
            frame.line.width = Pt(1.5)

        elif stijl == "playful":
            bg_shape.fill.fore_color.rgb = RGBColor(*config["bg"])
            if layout == "title":
                shadow = slide.shapes.add_shape(1, Inches(1.6), Inches(1.3), Inches(7.0), Inches(3.0))
                shadow.fill.solid()
                shadow.fill.fore_color.rgb = RGBColor(40,40,40)
                shadow.line.fill.background()
                vlak = slide.shapes.add_shape(1, Inches(1.5), Inches(1.2), Inches(7.0), Inches(3.0))
                vlak.fill.solid()
                vlak.fill.fore_color.rgb = RGBColor(*config["accent"])
                vlak.line.color.rgb = RGBColor(40,40,40)
                vlak.line.width = Pt(3)
            else:
                t_vlak = slide.shapes.add_shape(1, Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
                t_vlak.fill.solid()
                t_vlak.fill.fore_color.rgb = RGBColor(*config["accent"])
                t_vlak.line.color.rgb = RGBColor(40, 40, 40)
                t_vlak.line.width = Pt(2.5)
                
                b_width = Inches(4.5) if has_image else Inches(8.8)
                b_vlak = slide.shapes.add_shape(1, Inches(0.6), Inches(1.7), b_width, Inches(3.3))
                b_vlak.fill.solid()
                b_vlak.fill.fore_color.rgb = RGBColor(*config["accent2"])
                b_vlak.line.color.rgb = RGBColor(40, 40, 40)
                b_vlak.line.width = Pt(2)

        elif stijl == "professional":
            if layout == "title":
                bg_shape.fill.fore_color.rgb = RGBColor(*config["accent"])
                line = slide.shapes.add_shape(1, Inches(3), Inches(2.7), Inches(4), Inches(0.05))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(*config["accent2"])
                line.line.fill.background()
            else:
                bg_shape.fill.fore_color.rgb = RGBColor(*config["bg"])
                header = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(1.4))
                header.fill.solid()
                header.fill.fore_color.rgb = RGBColor(*config["accent"])
                header.line.fill.background()

        elif stijl == "modern":
            bg_shape.fill.fore_color.rgb = RGBColor(*config["bg"])
            if layout == "title":
                circle = slide.shapes.add_shape(9, Inches(3), Inches(0.8), Inches(4), Inches(4))
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(*config["accent2"])
                circle.line.fill.background()
            else:
                if not has_image: 
                    circle = slide.shapes.add_shape(9, Inches(7.5), Inches(3), Inches(4, 4))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = RGBColor(*config["accent2"])
                    circle.line.fill.background()
                line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(3), Inches(0.04))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(*config["accent"])
                line.line.fill.background()

        # ==========================================
        # STAP 2: DOWNLOAD DE ECHTE UNSPLASH AFBEELDING
        # ==========================================
        if has_image:
            try:
                temp_img_path = os.path.join(output_dir, f"temp_img_{slide_info.get('nr', 1)}.jpg")
                
                # Download met browser-header simulatie tegen 403 blokkades
                req = urllib.request.Request(img_url, headers={'User-Agent': 'Mozilla/5.0'})
                with urllib.request.urlopen(req) as response, open(temp_img_path, 'wb') as out_file:
                    out_file.write(response.read())
                
                slide.shapes.add_picture(temp_img_path, Inches(5.5), Inches(1.6), width=Inches(3.8), height=Inches(3.3))
            except Exception as img_err:
                print(f"Waarschuwing: Kon Unsplash foto niet downloaden: {img_err}")

        # ==========================================
        # STAP 3: TEKST POSITIONERING
        # ==========================================
        if layout == "title":
            t_left, t_top, t_width, t_height = Inches(1.5), Inches(1.8), Inches(7.0), Inches(1.0)
            b_left, b_top, b_width, b_height = Inches(1.5), Inches(2.8), Inches(7.0), Inches(1.2)
            align = PP_ALIGN.CENTER
            t_color = RGBColor(255,255,255) if stijl in ["professional", "playful"] else RGBColor(*config["text"])
            b_color = RGBColor(255,255,255) if stijl == "professional" else RGBColor(40,40,40)
        else:
            t_left, t_top, t_width, t_height = Inches(0.8), Inches(0.4), Inches(8.5), Inches(1.0)
            b_width = Inches(4.3) if has_image else Inches(8.5)
            b_left, b_top, b_height = Inches(0.8), Inches(1.8), Inches(3.2)
            align = PP_ALIGN.LEFT
            t_color = RGBColor(255,255,255) if stijl in ["professional", "playful"] else RGBColor(*config["text"])
            b_color = RGBColor(40,40,40) if stijl in ["professional", "playful"] else RGBColor(*config["text"])

        t_box = slide.shapes.add_textbox(t_left, t_top, t_width, t_height)
        t_box.text_frame.text = slide_info.get('titel', 'Titel')
        t_box.text_frame.paragraphs[0].alignment = align
        t_box.text_frame.paragraphs[0].font.bold = True
        t_box.text_frame.paragraphs[0].font.size = Pt(38)
        t_box.text_frame.paragraphs[0].font.color.rgb = t_color
        t_box.text_frame.paragraphs[0].font.name = "Arial"

        b_box = slide.shapes.add_textbox(b_left, b_top, b_width, b_height)
        b_box.text_frame.word_wrap = True
        text = slide_info.get('body', '')
        paragraphs = text.split('\n')
        b_box.text_frame.text = paragraphs[0] if paragraphs else ''
        for p_text in paragraphs[1:]:
            p = b_box.text_frame.add_paragraph()
            p.text = p_text
        
        for p in b_box.text_frame.paragraphs:
            p.alignment = align
            p.font.size = Pt(20)
            p.font.color.rgb = b_color
            p.font.name = "Arial"

    def run():
        output_dir = sys.argv[1]
        data = json.loads(sys.stdin.read())
        slides = data.get('proposal', [])
        stijl = data.get('stijl', 'modern')
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
        for s in slides: create_slide(prs, s, stijl, output_dir)
        if not os.path.exists(output_dir): os.makedirs(output_dir)
        prs.save(os.path.join(output_dir, "presentatie.pptx"))
        print("presentatie.pptx")

    if __name__ == "__main__": run()
except Exception:
    print("FATALE FOUT IN PYTHON:\n" + traceback.format_exc())